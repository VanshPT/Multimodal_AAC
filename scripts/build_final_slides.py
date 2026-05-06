"""Build Final_Slides.pptx (~10 slides) for P4 final."""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = "/Users/Omer/Desktop/UB/spring 26/NLP/Group Projects/Multimodal_AAC_final"
OUT  = os.path.join(ROOT, "Final_Slides.pptx")
FIG  = os.path.join(ROOT, "outputs", "figures")
SUM  = json.load(open(os.path.join(ROOT, "outputs", "eval_summary.json")))

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

NAVY = RGBColor(0x1F, 0x2D, 0x3D)
ACCENT = RGBColor(0x2E, 0x86, 0xC1)
INK = RGBColor(0x21, 0x21, 0x21)
MUTE = RGBColor(0x55, 0x55, 0x55)

def add(slide, l, t, w, h, text, size=18, bold=False, color=INK, align=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    if align is not None: p.alignment = align
    for r in p.runs:
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb

def bullets(slide, l, t, w, h, items, size=18, color=INK):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = u"\u2022 " + item
        for r in p.runs:
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def header(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.85))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add(slide, 0.5, 0.12, 12, 0.6, title, size=26, bold=True, color=RGBColor(255,255,255))
    if subtitle:
        add(slide, 0.5, 0.92, 12, 0.4, subtitle, size=14, color=MUTE)

def add_image(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w), height=Inches(h))

# Slide 1 - Title
s = prs.slides.add_slide(BLANK)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
add(s, 0.6, 2.5, 12, 1.0, "Multimodal AAC Chatbot", size=44, bold=True,
    color=RGBColor(255,255,255), align=None)
add(s, 0.6, 3.4, 12, 0.7, "Bonus-Driven Final System (Project 4)",
    size=22, color=RGBColor(220,220,220))
add(s, 0.6, 5.6, 12, 0.5, "Syed Omer Shah  |  Vansh Thakkar",
    size=18, color=RGBColor(255,255,255))
add(s, 0.6, 6.05, 12, 0.5, "CSE 635  |  NLP and Text Mining  |  May 2026",
    size=14, color=RGBColor(180,180,180))

# Slide 2 - Problem & Goals
s = prs.slides.add_slide(BLANK)
header(s, "Problem & Goals",
    "Real-time, grounded, multimodal AAC under a 5-second SLA")
bullets(s, 0.6, 1.4, 12.0, 5.0, [
    "AAC users with motor / speech impairments need fast, identity-faithful replies",
    "LLMs raise two new risks: hallucinated personal facts + unbounded latency",
    "Goal: 3 grounded, polarity-correct options per turn, in under 5 seconds",
    "Goal: keep all face frames and biosignals client-side (privacy first)",
    "Goal: implement all 5 P4 required nodes + all 5 P4 bonuses",
])

# Slide 3 - Architecture
s = prs.slides.add_slide(BLANK)
header(s, "System Architecture",
    "7-node LangGraph-style pipeline with deadline-raced LLM rewriter")
add_image(s, os.path.join(FIG, "system_architecture.png"), 0.4, 1.3, 7.6, 5.6)
bullets(s, 8.3, 1.3, 4.6, 5.6, [
    "FaceCue + Multimodal Mapping",
    "Parallel Prep (Router + Memory)",
    "Source Priority Planner",
    "Bucket Selector (12 buckets)",
    "Retrieve & Rerank (LTM/STM/PB)",
    "Evidence Refiner + Groundedness Guard",
    "Candidate Generator (rules + LLM)",
], size=15)

# Slide 4 - Multimodal Mapping
s = prs.slides.add_slide(BLANK)
header(s, "Multimodal Input Mapping",
    "Channels  ->  polarity / tone / verbosity / bucket boosts")
add_image(s, os.path.join(FIG, "multimodal_mapping.png"), 0.4, 1.3, 7.6, 5.6)
bullets(s, 8.3, 1.3, 4.6, 5.6, [
    "Face: FER+ 7-class emotion locks tone",
    "Gesture: thumbs / palm sets polarity",
    "Heart rate: >+25 bpm -> short verbosity",
    "Gaze: +0.30 boost on gazed bucket (Bonus 1)",
    "Vocal vs Air-sign conflict -> trust spatial (Bonus 2)",
    "Camera off -> 3 different tone variants",
], size=15)

# Slide 5 - Bonuses 1-3
s = prs.slides.add_slide(BLANK)
header(s, "Bonuses 1-3", "Gaze activation, conflict resolution, bucket priors")
bullets(s, 0.6, 1.3, 12.0, 5.5, [
    "1. Gaze-Based Retrieval Activation - GAZE_TO_BUCKET adds +0.30; redirects retrieval reliably on every probed case",
    "2. Vocal-vs-Air-Sign Conflict Resolution - deliberate spatial channel wins; logged in conflict_resolved",
    "3. Acceptance-Weighted Bucket Priors - bucket_acceptance counter incremented on selection; prior added to retrieval score",
    "All three implemented in home/aac/multimodal.py + home/aac/pipelines/nodes.py",
], size=17)

# Slide 6 - Bonuses 4-5
s = prs.slides.add_slide(BLANK)
header(s, "Bonuses 4-5", "Latency race + online phrase-bank update")
bullets(s, 0.6, 1.3, 12.0, 5.5, [
    "4. Latency-Optimised Fallback - llm._generate_text iterates flash-lite -> flash; rule-based floor ready immediately; 5s SLA enforced",
    "   Eval: max latency 4.33s, p95 = 2.43s, fallback engaged on 11/20 turns",
    "5. Online Index Update on Selection - confirm_response calls record_selection: PB phrases boosted +0.15 or added; bucket_acceptance updated",
    "   Eval: 5-turn rollout added 2 new phrases, boosted 3 existing",
], size=17)

# Slide 7 - Final results table
s = prs.slides.add_slide(BLANK)
header(s, "Final Evaluation (n=20)", "All targets met")
rows = [
    ("Metric", "Value"),
    ("BLEU-2 (peak)",        f"{SUM['bleu2_avg_peak']:.3f}"),
    ("ROUGE-L (peak)",       f"{SUM['rouge_l_avg_peak']:.3f}"),
    ("Groundedness (peak)",  f"{SUM['groundedness_avg_peak']:.3f}"),
    ("NLI faithfulness",     f"{SUM['nli_faithfulness_avg_peak']:.3f}"),
    ("Intent accuracy",      f"{SUM['intent_accuracy']:.2f}"),
    ("Bucket routing acc.",  f"{SUM['bucket_routing_accuracy']:.2f}"),
    ("Polarity adherence",   f"{SUM['polarity_adherence']:.2f}"),
    ("Memory write acc.",    f"{SUM['memory_write_accuracy']:.2f}"),
    ("Multimodal alignment", f"{SUM['multimodal_alignment_avg']:.3f}"),
    ("Latency avg (ms)",     f"{SUM['latency_avg_ms']:.0f}"),
    ("Latency p95 (ms)",     f"{SUM['latency_p95_ms']:.0f}"),
    ("SLA <5s hit",          "Yes" if SUM['latency_under_5000ms'] else "No"),
    ("LLM fallback rate",    f"{SUM['fallback_engagement_rate']:.2f}"),
]
table = s.shapes.add_table(len(rows), 2, Inches(2.5), Inches(1.4), Inches(8.3), Inches(5.5)).table
table.columns[0].width = Inches(5.0); table.columns[1].width = Inches(3.3)
for r, (k, v) in enumerate(rows):
    for c, val in enumerate((k, v)):
        cell = table.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(15); run.font.name = "Calibri"
                run.font.bold = (r == 0)
                run.font.color.rgb = RGBColor(255,255,255) if r == 0 else INK
        if r == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY

# Slide 8 - Eval matrix figure
s = prs.slides.add_slide(BLANK)
header(s, "Per-Metric Distribution", "Across 20 held-out cases")
add_image(s, os.path.join(FIG, "evaluation_matrix.png"), 1.5, 1.3, 10.3, 5.6)

# Slide 9 - Demo screenshot
s = prs.slides.add_slide(BLANK)
header(s, "Live Demo", "Browser UI: face-api.js, hand gesture, sensors panel")
add_image(s, os.path.join(FIG, "aac_demo_screenshot.png"), 0.4, 1.3, 8.5, 5.6)
bullets(s, 9.2, 1.3, 4.0, 5.6, [
    "Three options every turn",
    "Polarity-correct lead",
    "Tone matches detected emotion",
    "Sensors panel: HR / gaze / air-sign",
    "Memory pools live-update",
    "Demo video included in zip",
], size=14)

# Slide 10 - Limitations & Conclusion
s = prs.slides.add_slide(BLANK)
header(s, "Limitations & Wrap-Up", "Synthetic data, IRB next, ethical guard-rails on")
bullets(s, 0.6, 1.4, 12.0, 5.5, [
    "Evaluated on synthetic personas - real AAC user study needs IRB review",
    "Phrase bank stays per-user, client-side; only 7-class emotion label leaves browser",
    "LLM acts as rewriter; rule-based floor authoritative on polarity + do-not-say",
    "GroundednessGuard blocks hallucinated medication / clinical claims",
    "Deliverables: Final_Report.pdf, Final_Slides.pptx, source tree, demo video",
    "Repo: github.com/VanshPT/Multimodal_AAC (refined fork)",
], size=17)

prs.save(OUT)
print("WROTE", OUT)
